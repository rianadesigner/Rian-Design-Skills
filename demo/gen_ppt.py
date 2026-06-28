#!/usr/bin/env python3
"""Generate LLM Wiki architecture diagram PPT in Chinese - dual layer version."""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Colors
BG_COLOR = RGBColor(0x16, 0x16, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x80)
SECTION_LABEL = RGBColor(0x55, 0x55, 0x70)

BLUE_BORDER = RGBColor(0x44, 0x88, 0xFF)
GREEN_BORDER = RGBColor(0x33, 0xCC, 0x77)
PURPLE_BORDER = RGBColor(0x77, 0x55, 0xCC)
PINK_BORDER = RGBColor(0xCC, 0x33, 0x77)
ORANGE_BORDER = RGBColor(0xFF, 0x88, 0x33)
CYAN_BORDER = RGBColor(0x33, 0xCC, 0xCC)

BLUE_TEXT = RGBColor(0x88, 0xBB, 0xFF)
GREEN_TEXT = RGBColor(0x55, 0xEE, 0xBB)
PURPLE_TEXT = RGBColor(0xAA, 0x88, 0xFF)
PINK_TEXT = RGBColor(0xFF, 0x88, 0xBB)
YELLOW_TEXT = RGBColor(0xFF, 0xCC, 0x55)
LILAC_TEXT = RGBColor(0xCC, 0x88, 0xFF)
ORANGE_TEXT = RGBColor(0xFF, 0xAA, 0x55)
CYAN_TEXT = RGBColor(0x55, 0xEE, 0xEE)

TAG_BG = RGBColor(0x22, 0x22, 0x3A)
TAG_BORDER = RGBColor(0x3A, 0x3A, 0x55)
TAG_TEXT = RGBColor(0x99, 0x99, 0xBB)

blank_layout = prs.slide_layouts[6]


def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR


def add_text_box(slide, left, top, width, height, text, font_size=12,
                 color=WHITE, bold=False, alignment=PP_ALIGN.CENTER):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "PingFang SC"
    return txBox


def add_card(slide, left, top, width, height, border_color,
             title, title_color, subtitle, subtitle_color=GRAY_TEXT,
             fill_color=None):
    if fill_color is None:
        fill_color = RGBColor(0x1E, 0x2A, 0x3A)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    add_text_box(slide, left, top + Cm(0.6), width, Cm(1.2),
                 title, font_size=16, color=title_color, bold=True)
    add_text_box(slide, left, top + Cm(1.8), width, Cm(1),
                 subtitle, font_size=11, color=subtitle_color)


def add_tag(slide, left, top, text, w=Cm(3.2)):
    h = Cm(0.9)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h
    )
    shape.line.color.rgb = TAG_BORDER
    shape.line.width = Pt(1)
    shape.fill.solid()
    shape.fill.fore_color.rgb = TAG_BG
    tf = shape.text_frame
    tf.word_wrap = False
    tf.auto_size = None
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.color.rgb = TAG_TEXT
    run.font.name = "PingFang SC"


def add_dashed_line(slide, x1, y1, x2, y2, color=RGBColor(0x44, 0x44, 0x66)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(1.2)
    connector.line.dash_style = 4
    return connector


# ============================================================
# SLIDE 1: Original Architecture
# ============================================================
slide1 = prs.slides.add_slide(blank_layout)
set_bg(slide1)

add_text_box(slide1, Cm(0), Cm(0.5), Inches(16), Cm(2),
             "LLM 知识库", font_size=44, color=WHITE, bold=True)
add_text_box(slide1, Cm(0), Cm(2.5), Inches(16), Cm(1.2),
             "KARPATHY 知识编译模式", font_size=12, color=SECTION_LABEL)

# Architecture
add_text_box(slide1, Cm(0), Cm(4.0), Inches(16), Cm(0.8),
             "架  构", font_size=10, color=SECTION_LABEL)

arch_y = Cm(5.0)
card_h = Cm(3.0)

add_card(slide1, Cm(4), arch_y, Cm(7), card_h,
         BLUE_BORDER, "📁 原始资料", BLUE_TEXT, "不可变 · 只读",
         fill_color=RGBColor(0x1A, 0x2A, 0x40))
add_text_box(slide1, Cm(11), arch_y + Cm(0.8), Cm(3.5), Cm(1),
             "— 读取 →", font_size=11, color=GRAY_TEXT)
add_card(slide1, Cm(14.5), arch_y, Cm(7), card_h,
         GREEN_BORDER, "📒 知识库", GREEN_TEXT, "LLM 管理的 Markdown",
         fill_color=RGBColor(0x1A, 0x30, 0x2A))
add_text_box(slide1, Cm(21.5), arch_y + Cm(0.8), Cm(3.5), Cm(1),
             "← 指导 —", font_size=11, color=GRAY_TEXT)
add_card(slide1, Cm(25), arch_y, Cm(7.5), card_h,
         PURPLE_BORDER, "⚙️ 模式规范", PURPLE_TEXT, "CLAUDE.md / AGENTS.md",
         fill_color=RGBColor(0x22, 0x1A, 0x3A))

# Operations
add_text_box(slide1, Cm(0), Cm(9.2), Inches(16), Cm(0.8),
             "操  作", font_size=10, color=SECTION_LABEL)

ops_card_y = Cm(10.2)
ops_card_h = Cm(3.0)

add_card(slide1, Cm(4.5), ops_card_y, Cm(7), ops_card_h,
         PINK_BORDER, "📥 摄入", PINK_TEXT, "源文件 → 知识库更新",
         fill_color=RGBColor(0x2A, 0x1A, 0x2A))
add_card(slide1, Cm(14.5), ops_card_y, Cm(7), ops_card_h,
         PINK_BORDER, "🔍 查询", YELLOW_TEXT, "问题 → 综合回答",
         fill_color=RGBColor(0x2A, 0x1A, 0x2A))
add_card(slide1, Cm(24.5), ops_card_y, Cm(7), ops_card_h,
         PINK_BORDER, "🔧 检查", LILAC_TEXT, "健康检查 · 一致性",
         fill_color=RGBColor(0x2A, 0x1A, 0x2A))

# Steps
steps_y = Cm(13.8)
add_tag(slide1, Cm(4.5), steps_y, "读取源文件")
add_tag(slide1, Cm(8.0), steps_y, "写入摘要")
add_tag(slide1, Cm(5.2), steps_y + Cm(1.2), "更新索引")
add_tag(slide1, Cm(8.7), steps_y + Cm(1.2), "交叉链接")

add_tag(slide1, Cm(14.5), steps_y, "搜索索引")
add_tag(slide1, Cm(18.0), steps_y, "综合信息")
add_tag(slide1, Cm(15.5), steps_y + Cm(1.2), "归档回知识库")

add_tag(slide1, Cm(24.5), steps_y, "发现问题")
add_tag(slide1, Cm(28.0), steps_y, "修复补丁")
add_tag(slide1, Cm(25.8), steps_y + Cm(1.2), "建议来源")

# Actors
add_text_box(slide1, Cm(0), Cm(17.0), Inches(16), Cm(0.8),
             "参 与 者", font_size=10, color=SECTION_LABEL)

add_card(slide1, Cm(8), Cm(18.0), Cm(8), Cm(3.0),
         BLUE_BORDER, "🧑 人类", YELLOW_TEXT, "策划 · 提问 · 思考",
         fill_color=RGBColor(0x1A, 0x2A, 0x40))
add_card(slide1, Cm(20), Cm(18.0), Cm(8), Cm(3.0),
         GREEN_BORDER, "🤖 LLM 代理", GREEN_TEXT, "总结 · 交叉引用 · 维护",
         fill_color=RGBColor(0x1A, 0x30, 0x2A))

# Connections
add_dashed_line(slide1, Cm(8), Cm(13.2), Cm(12), Cm(18.0))
add_dashed_line(slide1, Cm(8), Cm(13.2), Cm(24), Cm(18.0))
add_dashed_line(slide1, Cm(18), Cm(13.2), Cm(12), Cm(18.0))
add_dashed_line(slide1, Cm(18), Cm(13.2), Cm(24), Cm(18.0))
add_dashed_line(slide1, Cm(28), Cm(13.2), Cm(24), Cm(18.0))


# ============================================================
# SLIDE 2: Dual-Layer Architecture
# ============================================================
slide2 = prs.slides.add_slide(blank_layout)
set_bg(slide2)

add_text_box(slide2, Cm(0), Cm(0.5), Inches(16), Cm(2),
             "LLM 知识库 · 双层架构", font_size=40, color=WHITE, bold=True)
add_text_box(slide2, Cm(0), Cm(2.5), Inches(16), Cm(1.2),
             "编译层（机器维护）  +  出版层（人类阅读）", font_size=13, color=SECTION_LABEL)

# ---- Left side: Compilation Layer ----
left_x = Cm(1.5)
layer_w = Cm(16)

add_text_box(slide2, left_x, Cm(4.0), layer_w, Cm(0.8),
             "编 译 层  ·  机 器 维 护", font_size=11, color=ORANGE_TEXT)

# Big container box for compilation layer
comp_box = slide2.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, left_x, Cm(4.8), layer_w, Cm(12.5)
)
comp_box.line.color.rgb = ORANGE_BORDER
comp_box.line.width = Pt(2)
comp_box.fill.solid()
comp_box.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x28)

# Compilation layer content
comp_inner_x = Cm(2.5)
comp_inner_y = Cm(5.5)

# Raw Sources
add_card(slide2, comp_inner_x, comp_inner_y, Cm(6), Cm(2.6),
         BLUE_BORDER, "📁 原始资料", BLUE_TEXT, "不可变 · 只读",
         fill_color=RGBColor(0x1A, 0x2A, 0x40))

# Arrow
add_text_box(slide2, comp_inner_x + Cm(6), comp_inner_y + Cm(0.6), Cm(2.5), Cm(1),
             "→", font_size=14, color=GRAY_TEXT)

# Wiki (Compilation)
add_card(slide2, comp_inner_x + Cm(7.5), comp_inner_y, Cm(7), Cm(2.6),
         GREEN_BORDER, "📒 编译知识库", GREEN_TEXT, "碎片化 · 强索引 · 交叉链接",
         fill_color=RGBColor(0x1A, 0x30, 0x2A))

# Characteristics tags
char_y = Cm(8.8)
add_tag(slide2, comp_inner_x, char_y, "小文件粒度", w=Cm(3.2))
add_tag(slide2, comp_inner_x + Cm(3.5), char_y, "向量嵌入", w=Cm(2.8))
add_tag(slide2, comp_inner_x + Cm(6.6), char_y, "Frontmatter", w=Cm(3.4))
add_tag(slide2, comp_inner_x + Cm(10.3), char_y, "Token 优化", w=Cm(3.0))

char_y2 = Cm(10.0)
add_tag(slide2, comp_inner_x, char_y2, "冗余但精确", w=Cm(3.4))
add_tag(slide2, comp_inner_x + Cm(3.7), char_y2, "自动更新", w=Cm(2.8))
add_tag(slide2, comp_inner_x + Cm(6.8), char_y2, "多粒度副本", w=Cm(3.4))

# Operations
ops_y2 = Cm(11.6)
add_card(slide2, comp_inner_x, ops_y2, Cm(4.2), Cm(2.2),
         PINK_BORDER, "📥 摄入", PINK_TEXT, "源 → 碎片",
         fill_color=RGBColor(0x2A, 0x1A, 0x2A))
add_card(slide2, comp_inner_x + Cm(4.8), ops_y2, Cm(4.2), Cm(2.2),
         PINK_BORDER, "🔍 查询", YELLOW_TEXT, "检索 → 注入",
         fill_color=RGBColor(0x2A, 0x1A, 0x2A))
add_card(slide2, comp_inner_x + Cm(9.6), ops_y2, Cm(4.2), Cm(2.2),
         PINK_BORDER, "🔧 检查", LILAC_TEXT, "一致性校验",
         fill_color=RGBColor(0x2A, 0x1A, 0x2A))

# Agent
add_card(slide2, comp_inner_x + Cm(3.5), Cm(14.4), Cm(7), Cm(2.4),
         GREEN_BORDER, "🤖 LLM 代理", GREEN_TEXT, "自动维护 · 索引 · 链接",
         fill_color=RGBColor(0x1A, 0x30, 0x2A))

# ---- Right side: Publication Layer ----
right_x = Cm(19.5)
pub_w = Cm(16)

add_text_box(slide2, right_x, Cm(4.0), pub_w, Cm(0.8),
             "出 版 层  ·  人 类 阅 读", font_size=11, color=CYAN_TEXT)

# Big container box for publication layer
pub_box = slide2.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, right_x, Cm(4.8), pub_w, Cm(12.5)
)
pub_box.line.color.rgb = CYAN_BORDER
pub_box.line.width = Pt(2)
pub_box.fill.solid()
pub_box.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x28)

pub_inner_x = Cm(20.5)
pub_inner_y = Cm(5.5)

# Publication Wiki
add_card(slide2, pub_inner_x, pub_inner_y, Cm(7), Cm(2.6),
         CYAN_BORDER, "📖 出版知识库", CYAN_TEXT, "叙事化 · 聚合 · 连贯",
         fill_color=RGBColor(0x1A, 0x2A, 0x35))

# Output formats
add_card(slide2, pub_inner_x + Cm(8), pub_inner_y, Cm(5.5), Cm(2.6),
         PURPLE_BORDER, "📤 输出格式", PURPLE_TEXT, "文档 · 报告 · 图表",
         fill_color=RGBColor(0x22, 0x1A, 0x3A))

# Characteristics
pub_char_y = Cm(8.8)
add_tag(slide2, pub_inner_x, pub_char_y, "按主题聚合", w=Cm(3.4))
add_tag(slide2, pub_inner_x + Cm(3.7), pub_char_y, "有上下文", w=Cm(2.8))
add_tag(slide2, pub_inner_x + Cm(6.8), pub_char_y, "去重精简", w=Cm(2.8))
add_tag(slide2, pub_inner_x + Cm(10.0), pub_char_y, "可读优先", w=Cm(2.8))

pub_char_y2 = Cm(10.0)
add_tag(slide2, pub_inner_x, pub_char_y2, "目录导航", w=Cm(2.8))
add_tag(slide2, pub_inner_x + Cm(3.1), pub_char_y2, "摘要总结", w=Cm(2.8))
add_tag(slide2, pub_inner_x + Cm(6.2), pub_char_y2, "可视化图表", w=Cm(3.4))

# Operations
pub_ops_y = Cm(11.6)
add_card(slide2, pub_inner_x, pub_ops_y, Cm(4.2), Cm(2.2),
         PINK_BORDER, "✍️ 编辑", PINK_TEXT, "整合 → 叙事",
         fill_color=RGBColor(0x2A, 0x1A, 0x2A))
add_card(slide2, pub_inner_x + Cm(4.8), pub_ops_y, Cm(4.5), Cm(2.2),
         PINK_BORDER, "📊 可视化", YELLOW_TEXT, "数据 → 图表",
         fill_color=RGBColor(0x2A, 0x1A, 0x2A))
add_card(slide2, pub_inner_x + Cm(9.8), pub_ops_y, Cm(4), Cm(2.2),
         PINK_BORDER, "✅ 审核", LILAC_TEXT, "质量把关",
         fill_color=RGBColor(0x2A, 0x1A, 0x2A))

# Human
add_card(slide2, pub_inner_x + Cm(3.5), Cm(14.4), Cm(7), Cm(2.4),
         BLUE_BORDER, "🧑 人类", YELLOW_TEXT, "策划 · 审核 · 标注质量",
         fill_color=RGBColor(0x1A, 0x2A, 0x40))

# ---- Connection between two layers (center arrow) ----
# Bidirectional arrow between layers
mid_x = Cm(17.8)
add_text_box(slide2, mid_x - Cm(0.5), Cm(6.0), Cm(2.5), Cm(1),
             "→", font_size=20, color=RGBColor(0x55, 0x55, 0x77))
add_text_box(slide2, mid_x - Cm(0.5), Cm(6.8), Cm(2.5), Cm(0.8),
             "素材供给", font_size=9, color=GRAY_TEXT)

add_text_box(slide2, mid_x - Cm(0.5), Cm(14.8), Cm(2.5), Cm(1),
             "←", font_size=20, color=RGBColor(0x55, 0x55, 0x77))
add_text_box(slide2, mid_x - Cm(0.5), Cm(15.6), Cm(2.5), Cm(0.8),
             "质量信号", font_size=9, color=GRAY_TEXT)

# Vertical dashed lines connecting layers
add_dashed_line(slide2, Cm(18.2), Cm(6.5), Cm(18.2), Cm(7.5),
                color=RGBColor(0x55, 0x55, 0x77))
add_dashed_line(slide2, Cm(18.2), Cm(14.5), Cm(18.2), Cm(15.5),
                color=RGBColor(0x55, 0x55, 0x77))


# ============================================================
# SLIDE 3: Data Flow between layers
# ============================================================
slide3 = prs.slides.add_slide(blank_layout)
set_bg(slide3)

add_text_box(slide3, Cm(0), Cm(0.5), Inches(16), Cm(2),
             "双层数据流", font_size=40, color=WHITE, bold=True)
add_text_box(slide3, Cm(0), Cm(2.5), Inches(16), Cm(1.2),
             "编译层与出版层之间的协作循环", font_size=13, color=SECTION_LABEL)

# Central flow diagram
flow_center = Cm(18)  # center x of slide

# Left: Compilation layer box
add_text_box(slide3, Cm(2), Cm(4.5), Cm(14), Cm(1),
             "编译层（LLM 维护）", font_size=14, color=ORANGE_TEXT, bold=True)

comp_flow_y = Cm(5.8)
add_card(slide3, Cm(2), comp_flow_y, Cm(5.5), Cm(2.5),
         BLUE_BORDER, "📁 原始资料", BLUE_TEXT, "代码 · 文档 · API",
         fill_color=RGBColor(0x1A, 0x2A, 0x40))
add_text_box(slide3, Cm(7.5), comp_flow_y + Cm(0.6), Cm(1.5), Cm(1),
             "→", font_size=16, color=GRAY_TEXT)
add_card(slide3, Cm(9), comp_flow_y, Cm(6), Cm(2.5),
         GREEN_BORDER, "🤖 LLM 处理", GREEN_TEXT, "摄入 · 索引 · 链接",
         fill_color=RGBColor(0x1A, 0x30, 0x2A))

# Right: Publication layer box
add_text_box(slide3, Cm(21), Cm(4.5), Cm(14), Cm(1),
             "出版层（人类消费）", font_size=14, color=CYAN_TEXT, bold=True)

add_card(slide3, Cm(21), comp_flow_y, Cm(6), Cm(2.5),
         CYAN_BORDER, "📖 知识出版", CYAN_TEXT, "叙事 · 图表 · 导航",
         fill_color=RGBColor(0x1A, 0x2A, 0x35))
add_text_box(slide3, Cm(27), comp_flow_y + Cm(0.6), Cm(1.5), Cm(1),
             "→", font_size=16, color=GRAY_TEXT)
add_card(slide3, Cm(28.5), comp_flow_y, Cm(6), Cm(2.5),
         PURPLE_BORDER, "🧑 人类阅读", YELLOW_TEXT, "理解 · 决策 · 反馈",
         fill_color=RGBColor(0x22, 0x1A, 0x3A))

# Forward flow arrow (compilation → publication)
arrow_y1 = Cm(9.5)
add_text_box(slide3, Cm(14), arrow_y1, Cm(9), Cm(1.5),
             "─── 素材供给 ──────→", font_size=13, color=ORANGE_TEXT)

# Description cards for forward flow
fwd_y = Cm(11.0)
add_tag(slide3, Cm(3), fwd_y, "碎片化知识", w=Cm(3.4))
add_tag(slide3, Cm(6.8), fwd_y, "结构化数据", w=Cm(3.4))
add_tag(slide3, Cm(10.6), fwd_y, "交叉引用图", w=Cm(3.4))

add_text_box(slide3, Cm(14.5), fwd_y - Cm(0.1), Cm(2), Cm(1),
             "→", font_size=16, color=ORANGE_TEXT)

add_tag(slide3, Cm(17), fwd_y, "聚合成文章", w=Cm(3.4))
add_tag(slide3, Cm(20.8), fwd_y, "生成图表", w=Cm(2.8))
add_tag(slide3, Cm(24.0), fwd_y, "构建目录", w=Cm(2.8))

# Backward flow arrow (publication → compilation)
arrow_y2 = Cm(13.5)
add_text_box(slide3, Cm(14), arrow_y2, Cm(9), Cm(1.5),
             "←────── 质量信号 ───", font_size=13, color=CYAN_TEXT)

# Description cards for backward flow
bwd_y = Cm(15.2)
add_tag(slide3, Cm(17), bwd_y, "标注错误", w=Cm(2.8))
add_tag(slide3, Cm(20.2), bwd_y, "补充来源", w=Cm(2.8))
add_tag(slide3, Cm(23.4), bwd_y, "优先级排序", w=Cm(3.4))

add_text_box(slide3, Cm(14.5), bwd_y - Cm(0.1), Cm(2), Cm(1),
             "←", font_size=16, color=CYAN_TEXT)

add_tag(slide3, Cm(3), bwd_y, "修正知识图", w=Cm(3.4))
add_tag(slide3, Cm(6.8), bwd_y, "重新摄入", w=Cm(2.8))
add_tag(slide3, Cm(10.0), bwd_y, "更新索引", w=Cm(2.8))

# Bottom summary
summary_y = Cm(18.0)
add_text_box(slide3, Cm(2), summary_y, Inches(16) - Cm(4), Cm(2),
             "编译层为出版层提供素材 ← → 出版层为编译层提供质量信号\n形成持续改进的闭环",
             font_size=14, color=RGBColor(0x88, 0x88, 0xAA))


# Save
output_path = "/Users/rian/Rian-Design-Skills/Demo/LLM知识库-Karpathy知识编译模式.pptx"
prs.save(output_path)
print(f"PPT saved: {output_path}")
